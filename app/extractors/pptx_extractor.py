import io
import pptx
from typing import List
from app.extractors.base import BaseExtractor, format_markdown_table, format_image_placeholder
from app.models import ExtractionResult, ExtractedImage
from app.assets import default_asset_store
from app.vision import vision_client


class PPTXExtractor(BaseExtractor):
    async def extract(self, file_bytes: bytes, filename: str, **kwargs) -> ExtractionResult:
        prs = pptx.Presentation(io.BytesIO(file_bytes))

        extracted_images: List[ExtractedImage] = []
        content_parts = []
        image_count = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_header = f"## Slide {slide_idx}"
            slide_content = [slide_header]

            for shape in slide.shapes:
                # Text frame
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_content.append(text)

                # Tables
                elif shape.has_table:
                    table = shape.table
                    rows_data = []
                    for row in table.rows:
                        row_cells = [cell.text.strip().replace("\n", "<br>") for cell in row.cells]
                        rows_data.append(row_cells)
                    if rows_data:
                        headers = rows_data[0]
                        rows = rows_data[1:] if len(rows_data) > 1 else []
                        slide_content.append(format_markdown_table(headers, rows))

                # Images
                elif shape.shape_type == pptx.enum.shapes.MSO_SHAPE.PICTURE or hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext or "png"
                        image_count += 1
                        img_name = f"{filename}_slide{slide_idx}_img_{image_count}.{image_ext}"
                        location = f"Slide {slide_idx}"
                        alt_text = f"Diagram / Image on Slide {slide_idx}"
                        asset_id, url = default_asset_store.save(
                            img_name, f"image/{image_ext}", alt_text, location, image_bytes
                        )
                        ext_img = ExtractedImage(
                            id=asset_id,
                            filename=img_name,
                            content_type=f"image/{image_ext}",
                            size_bytes=len(image_bytes),
                            width=shape.width or 0,
                            height=shape.height or 0,
                            alt_text=alt_text,
                            location=location,
                            relative_path=f"./assets/{img_name}",
                            url=url,
                            data=image_bytes,
                        )
                        extracted_images.append(ext_img)
                        slide_content.append(format_image_placeholder(ext_img))
                    except Exception:
                        pass

            content_parts.append("\n\n".join(slide_content))

        if vision_client and vision_client.enabled and extracted_images:
            await vision_client.enrich_images(extracted_images)

        final_content = "\n\n---\n\n".join(content_parts)
        word_count = len(final_content.split())

        return ExtractionResult(
            content=final_content,
            images=extracted_images,
            metadata={
                "filename": filename,
                "total_slides": len(prs.slides),
                "engine": "local",
            },
            word_count=word_count,
            detected_type="pptx",
        )
